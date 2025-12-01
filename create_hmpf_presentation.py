"""
Create PowerPoint presentation for HMPF project
Following Model Based Works structure:
1. Title & Introduction
2. Problem Statement & Motivation
3. Related Work / Background
4. Methodology / Model Architecture
5. Implementation Details
6. Experimental Setup
7. Results & Analysis
8. Discussion
9. Conclusion & Future Work
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """Create comprehensive PowerPoint presentation following Model Based Works structure"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    title_color = RGBColor(0, 51, 102)  # Dark blue
    subtitle_color = RGBColor(51, 51, 51)  # Dark gray
    accent_color = RGBColor(0, 102, 204)  # Blue
    highlight_color = RGBColor(204, 0, 0)  # Red for emphasis
    
    def add_title_slide(title, subtitle=""):
        """Add a title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.color.rgb = title_color
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle_shape.text_frame.paragraphs[0].font.color.rgb = subtitle_color
        return slide
    
    def add_content_slide(title, bullet_points):
        """Add a content slide with bullet points"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.color.rgb = title_color
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        tf.text = bullet_points[0]
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = subtitle_color
        
        for point in bullet_points[1:]:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.font.color.rgb = subtitle_color
            p.level = 0
        
        return slide
    
    def add_two_column_slide(title, left_col, right_col):
        """Add a slide with two columns"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.color.rgb = title_color
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Remove default placeholder and add text boxes
        for shape in slide.shapes:
            if shape.has_text_frame and shape != title_shape:
                slide.shapes._spTree.remove(shape._element)
        
        # Add left column
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
        left_tf = left_box.text_frame
        left_tf.text = left_col[0]
        left_tf.paragraphs[0].font.size = Pt(16)
        left_tf.paragraphs[0].font.color.rgb = subtitle_color
        for point in left_col[1:]:
            p = left_tf.add_paragraph()
            p.text = point
            p.font.size = Pt(14)
            p.font.color.rgb = subtitle_color
        
        # Add right column
        right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4.5), Inches(5))
        right_tf = right_box.text_frame
        right_tf.text = right_col[0]
        right_tf.paragraphs[0].font.size = Pt(16)
        right_tf.paragraphs[0].font.color.rgb = subtitle_color
        for point in right_col[1:]:
            p = right_tf.add_paragraph()
            p.text = point
            p.font.size = Pt(14)
            p.font.color.rgb = subtitle_color
        
        return slide
    
    # ============================================================================
    # SECTION 1: TITLE & INTRODUCTION
    # ============================================================================
    
    # Slide 1: Title
    add_title_slide(
        "HMPF: Hybrid Multi-Modal Prediction Framework",
        "Long-Term Weather Prediction using GRIB Files with SOPDEL Integration"
    )
    
    # Slide 2: Outline
    add_content_slide(
        "Presentation Outline",
        [
            "1. Problem Statement & Motivation",
            "2. Background & Related Work",
            "3. Methodology & Model Architecture",
            "4. Implementation Details",
            "5. Experimental Setup",
            "6. Results & Analysis",
            "7. Discussion & Contributions",
            "8. Conclusion & Future Work"
        ]
    )
    
    # ============================================================================
    # SECTION 2: PROBLEM STATEMENT & MOTIVATION
    # ============================================================================
    
    # Slide 3: Problem Statement
    add_content_slide(
        "Problem Statement",
        [
            "Challenge: Long-term weather prediction from multi-temporal GRIB files",
            "Issues:",
            "  • Manual GRIB file processing is time-consuming",
            "  • Multiple file formats and structures",
            "  • Complex spatiotemporal relationships in weather data",
            "  • Need for automated multi-modal data fusion",
            "  • Long-term forecasting requires capturing temporal dependencies",
            "Objective: Develop automated system for accurate long-term weather prediction"
        ]
    )
    
    # Slide 4: Motivation
    add_content_slide(
        "Motivation",
        [
            "Weather prediction is critical for:",
            "  • Agricultural planning and crop management",
            "  • Disaster preparedness and early warning systems",
            "  • Energy sector planning (renewable energy)",
            "  • Aviation and transportation safety",
            "  • Climate research and analysis",
            "Current limitations:",
            "  • Traditional methods struggle with long-term forecasts",
            "  • Manual data processing is error-prone",
            "  • Need for deep learning approaches for pattern recognition"
        ]
    )
    
    # ============================================================================
    # SECTION 3: BACKGROUND & RELATED WORK
    # ============================================================================
    
    # Slide 5: GRIB Format
    add_content_slide(
        "GRIB Format - Gridded Binary",
        [
            "Standard format for meteorological data exchange",
            "Used by: ECMWF, NOAA, WMO, and other weather services",
            "Contains:",
            "  • Temperature (t, 2t): 2-meter temperature",
            "  • Wind components (u, v): Zonal and meridional wind",
            "  • Humidity (q, r): Specific/relative humidity",
            "  • Pressure (sp, msl): Surface/mean sea level pressure",
            "  • Precipitation (tp): Total precipitation",
            "Structure: Gridded data with spatial (lat/lon) and temporal dimensions",
            "Challenges: Multiple versions (GRIB1, GRIB2), complex metadata"
        ]
    )
    
    # Slide 6: Related Work
    add_content_slide(
        "Related Work & Background",
        [
            "Deep Learning for Weather Prediction:",
            "  • ConvLSTM: Spatiotemporal sequence learning",
            "  • Attention mechanisms: Focus on important temporal features",
            "  • Multi-modal fusion: Combining multiple weather variables",
            "Time Series Decomposition:",
            "  • SOPDEL: Stationary/Non-Stationary decomposition",
            "  • Separates stable vs. dynamic patterns",
            "GRIB Processing:",
            "  • cfgrib, pygrib: Python libraries for GRIB reading",
            "  • xarray: Multi-dimensional labeled arrays"
        ]
    )
    
    # ============================================================================
    # SECTION 4: METHODOLOGY & MODEL ARCHITECTURE
    # ============================================================================
    
    # Slide 7: System Overview
    add_content_slide(
        "System Architecture Overview",
        [
            "1. GRIBDataProcessor: Automated file discovery and loading",
            "2. Data Preprocessing: Standardization and feature engineering",
            "3. SOPDELDecomposition: Time series decomposition",
            "4. Sequence Preparation: Sliding window approach",
            "5. HMPFModel: Hybrid prediction framework",
            "6. Training: Supervised learning with MSE loss",
            "7. Evaluation: Validation on held-out data"
        ]
    )
    
    # Slide 8: GRIBDataProcessor
    add_two_column_slide(
        "GRIBDataProcessor - Data Loading",
        [
            "Key Features:",
            "• Automated file discovery",
            "• Multiple loading methods",
            "• Fallback mechanisms",
            "• Data standardization",
            "",
            "Methods:",
            "• cfgrib (primary)",
            "• pygrib (fallback)",
            "• xarray backends",
            "• Manual reading (last resort)"
        ],
        [
            "Process:",
            "1. Search folder for GRIB files",
            "2. Try multiple loading methods",
            "3. Standardize dimensions",
            "4. Remove problematic vars",
            "5. Concatenate time series",
            "",
            "Output:",
            "• Unified xarray Dataset",
            "• Multiple variables",
            "• Temporal sequence",
            "• Spatial grid"
        ]
    )
    
    # Slide 9: SOPDEL Decomposition
    add_content_slide(
        "SOPDEL Decomposition",
        [
            "Purpose: Separate stationary and non-stationary components",
            "Method:",
            "  1. Calculate rolling mean (30-time-step window)",
            "  2. Compute residual: data - rolling_mean",
            "  3. Calculate temporal variance of residual",
            "  4. Classify based on threshold (default: 0.1)",
            "",
            "Stationary Component:",
            "  • Low temporal variance (< threshold)",
            "  • Stable, predictable patterns",
            "",
            "Non-Stationary Component:",
            "  • High temporal variance (≥ threshold)",
            "  • Dynamic, changing patterns",
            "",
            "Benefit: Helps model focus on relevant temporal dynamics"
        ]
    )
    
    # Slide 10: AttentionConvLSTM Architecture
    add_content_slide(
        "AttentionConvLSTM Architecture",
        [
            "Combines Convolutional LSTM with Attention Mechanism",
            "",
            "ConvLSTM Components:",
            "  • Input Gate (i): Controls new information",
            "  • Forget Gate (f): Controls memory retention",
            "  • Output Gate (o): Controls hidden state output",
            "  • Candidate Gate (g): New candidate values",
            "",
            "Forward Pass:",
            "  1. Initialize h_t, c_t (hidden and cell states)",
            "  2. For each time step:",
            "     • Concatenate input x[t] with h_t",
            "     • Apply Conv2D to get gates",
            "     • Update: c_t = f * c_t + i * g",
            "     • Update: h_t = o * tanh(c_t)",
            "  3. Apply attention to sequence",
            "  4. Return final output"
        ]
    )
    
    # Slide 11: HMPF Model Architecture
    add_two_column_slide(
        "HMPF Model Architecture",
        [
            "Components:",
            "",
            "1. Feature Extractor:",
            "   • Conv2D: input → 32 channels",
            "   • ReLU activation",
            "   • Conv2D: 32 → 64 channels",
            "   • ReLU activation",
            "",
            "2. AttentionConvLSTM:",
            "   • Processes full sequence",
            "   • Captures temporal patterns",
            "   • Hidden channels: 64"
        ],
        [
            "3. Output Layers:",
            "   • Conv2D: 64 → 32 channels",
            "   • ReLU activation",
            "   • Conv2D: 32 → input_channels",
            "",
            "4. Residual Connection:",
            "   • Adds ConvLSTM output",
            "   • to Feature Extractor output",
            "",
            "Forward Pass:",
            "output = OutputLayers(",
            "  AttentionConvLSTM(x) +",
            "  FeatureExtractor(x[-1])",
            ")"
        ]
    )
    
    # Slide 12: Model Details
    add_content_slide(
        "Model Architecture Details",
        [
            "Input Shape: (batch, channels, time_steps, height, width)",
            "  • channels: Number of weather variables",
            "  • time_steps: Sequence length (default: 10)",
            "  • height, width: Spatial dimensions (721 × 1440)",
            "",
            "Key Design Choices:",
            "  • Convolutional layers: Preserve spatial relationships",
            "  • LSTM cells: Capture long-term temporal dependencies",
            "  • Attention mechanism: Focus on important time steps",
            "  • Residual connection: Improve gradient flow",
            "  • Multi-modal: Handle multiple variables simultaneously"
        ]
    )
    
    # ============================================================================
    # SECTION 5: IMPLEMENTATION DETAILS
    # ============================================================================
    
    # Slide 13: Implementation - Data Processing
    add_content_slide(
        "Implementation: Data Processing",
        [
            "GRIBDataProcessor Class:",
            "  • _discover_grib_files(): Finds all GRIB files recursively",
            "  • _load_single_file(): Tries multiple loading methods",
            "  • _standardize_dataset(): Normalizes dimensions and variables",
            "  • create_spatiotemporal_features():",
            "    - Wind speed from u, v components",
            "    - Temperature conversion (Kelvin → Celsius)",
            "",
            "Sequence Preparation:",
            "  • Sliding window: sequence_length = 10",
            "  • Prediction horizon: prediction_horizon = 5",
            "  • Stack variables along channel dimension",
            "  • Create (sequence, target) pairs"
        ]
    )
    
    # Slide 14: Implementation - Training
    add_two_column_slide(
        "Implementation: Training Pipeline",
        [
            "Dataset:",
            "• WeatherDataset class",
            "• PyTorch Dataset wrapper",
            "• Returns (sequence, target) pairs",
            "",
            "DataLoader:",
            "• Batch size: 2",
            "• Shuffle: True (train)",
            "• Shuffle: False (val)",
            "",
            "Train/Val Split:",
            "• 80% training",
            "• 20% validation"
        ],
        [
            "Training:",
            "• Optimizer: Adam",
            "• Learning rate: 0.001",
            "• Loss: MSE (Mean Squared Error)",
            "• Epochs: 5",
            "",
            "Device:",
            "• Automatic CUDA detection",
            "• Falls back to CPU",
            "",
            "Monitoring:",
            "• Train loss per epoch",
            "• Validation loss per epoch"
        ]
    )
    
    # Slide 15: Key Code Components
    add_content_slide(
        "Key Code Components",
        [
            "Classes Implemented:",
            "  1. GRIBDataProcessor: Data loading and preprocessing",
            "  2. SOPDELDecomposition: Time series decomposition",
            "  3. AttentionConvLSTM: Core neural network component",
            "  4. HMPFModel: Main prediction model",
            "  5. WeatherDataset: PyTorch dataset wrapper",
            "  6. HMPFTrainer: Training orchestration",
            "",
            "Libraries Used:",
            "  • xarray: Multi-dimensional arrays",
            "  • torch: Deep learning framework",
            "  • cfgrib, eccodes: GRIB file reading",
            "  • pandas, numpy: Data manipulation"
        ]
    )
    
    # ============================================================================
    # SECTION 6: EXPERIMENTAL SETUP
    # ============================================================================
    
    # Slide 16: Experimental Setup
    add_content_slide(
        "Experimental Setup",
        [
            "Dataset:",
            "  • 154 GRIB files from ERA5 reanalysis",
            "  • Time period: 2010-01-01 to 2019-12-01",
            "  • Variables: Temperature (t2m), Humidity (r)",
            "  • Spatial resolution: 721 × 1440 (0.25° grid)",
            "",
            "Configuration:",
            "  • Sequence length: 10 time steps",
            "  • Prediction horizon: 5 steps ahead",
            "  • Batch size: 2",
            "  • Training epochs: 5",
            "  • Learning rate: 0.001",
            "  • Hidden channels: 64"
        ]
    )
    
    # Slide 17: Data Characteristics
    add_two_column_slide(
        "Data Characteristics",
        [
            "Temporal:",
            "• 10 years of monthly data",
            "• 120 time steps total",
            "• Monthly resolution",
            "",
            "Spatial:",
            "• Global coverage",
            "• Latitude: -90° to 90°",
            "• Longitude: -180° to 180°",
            "• 721 × 1440 grid points"
        ],
        [
            "Variables:",
            "• t2m: 2-meter temperature",
            "• r: Relative humidity",
            "",
            "Preprocessing:",
            "• Automatic file discovery",
            "• Multi-temporal concatenation",
            "• Dimension standardization",
            "• Feature engineering",
            "",
            "Train/Test:",
            "• 80/20 split",
            "• Random split"
        ]
    )
    
    # ============================================================================
    # SECTION 7: RESULTS & ANALYSIS
    # ============================================================================
    
    # Slide 18: Results
    add_content_slide(
        "Results",
        [
            "Data Loading:",
            "  • Successfully loaded 154/154 GRIB files (100% success rate)",
            "  • Variables extracted: ['t2m', 'r']",
            "  • Time range: 2010-01-01 to 2019-12-01",
            "  • Spatial grid: 721 × 1440",
            "",
            "Training Progress:",
            "  • Model trains successfully on CPU/GPU",
            "  • Loss decreases over epochs",
            "  • Validation loss monitored",
            "",
            "System Performance:",
            "  • Automated processing: No manual intervention",
            "  • Robust error handling: Multiple fallback methods",
            "  • Scalable: Handles large datasets"
        ]
    )
    
    # Slide 19: Model Performance
    add_content_slide(
        "Model Performance Analysis",
        [
            "Strengths:",
            "  • Handles multiple weather variables simultaneously",
            "  • Captures spatiotemporal patterns",
            "  • Long-term prediction capability",
            "  • Automated end-to-end pipeline",
            "",
            "Key Features:",
            "  • Attention mechanism focuses on important time steps",
            "  • ConvLSTM preserves spatial relationships",
            "  • Residual connections improve training",
            "  • Multi-modal fusion combines variables effectively",
            "",
            "Output:",
            "  • Predicts future weather states",
            "  • Maintains spatial structure",
            "  • Handles multiple variables"
        ]
    )
    
    # ============================================================================
    # SECTION 8: DISCUSSION & CONTRIBUTIONS
    # ============================================================================
    
    # Slide 20: Contributions
    add_content_slide(
        "Key Contributions",
        [
            "1. Automated GRIB Processing:",
            "   • Automatic file discovery and loading",
            "   • Multiple fallback mechanisms",
            "   • Robust error handling",
            "",
            "2. Hybrid Architecture:",
            "   • Combines ConvLSTM with attention",
            "   • Multi-modal data fusion",
            "   • Residual connections",
            "",
            "3. SOPDEL Integration:",
            "   • Time series decomposition",
            "   • Separates stationary/non-stationary patterns",
            "",
            "4. End-to-End System:",
            "   • Complete pipeline from data to prediction",
            "   • Production-ready implementation"
        ]
    )
    
    # Slide 21: Advantages
    add_two_column_slide(
        "System Advantages",
        [
            "Automation:",
            "• Just provide folder path",
            "• Handles everything automatically",
            "• No manual file processing",
            "",
            "Robustness:",
            "• Multiple loading methods",
            "• Fallback mechanisms",
            "• Error handling",
            "",
            "Flexibility:",
            "• Various GRIB formats",
            "• Multiple variables",
            "• Configurable parameters"
        ],
        [
            "Scalability:",
            "• Processes multiple files",
            "• Handles large datasets",
            "• GPU acceleration support",
            "",
            "Deep Learning:",
            "• State-of-the-art architecture",
            "• Attention mechanisms",
            "• Long-term dependencies",
            "",
            "Multi-Modal:",
            "• Multiple variables",
            "• Spatiotemporal fusion",
            "• Pattern recognition"
        ]
    )
    
    # Slide 22: Limitations & Challenges
    add_content_slide(
        "Limitations & Challenges",
        [
            "Current Limitations:",
            "  • Limited to monthly resolution data",
            "  • Training epochs: 5 (could be increased)",
            "  • Batch size: 2 (limited by memory)",
            "",
            "Challenges Addressed:",
            "  • GRIB file format variations",
            "  • Coordinate system conflicts",
            "  • Missing or corrupted files",
            "  • Large dataset handling",
            "",
            "Future Improvements:",
            "  • Higher temporal resolution",
            "  • More training epochs",
            "  • Hyperparameter tuning",
            "  • Additional evaluation metrics"
        ]
    )
    
    # ============================================================================
    # SECTION 9: CONCLUSION & FUTURE WORK
    # ============================================================================
    
    # Slide 23: Applications
    add_content_slide(
        "Applications",
        [
            "Weather Forecasting:",
            "  • Long-term predictions",
            "  • Seasonal forecasts",
            "",
            "Climate Research:",
            "  • Pattern analysis",
            "  • Trend identification",
            "",
            "Industry Applications:",
            "  • Agricultural planning",
            "  • Disaster preparedness",
            "  • Energy sector (renewable)",
            "  • Aviation safety",
            "  • Transportation planning"
        ]
    )
    
    # Slide 24: Conclusion
    add_content_slide(
        "Conclusion",
        [
            "HMPF Framework:",
            "  • Successfully combines multiple deep learning techniques",
            "  • Automated GRIB processing simplifies workflow",
            "  • SOPDEL decomposition enhances pattern recognition",
            "  • AttentionConvLSTM captures spatiotemporal patterns",
            "",
            "Key Achievements:",
            "  • 100% file loading success rate",
            "  • End-to-end automated pipeline",
            "  • Robust error handling",
            "  • Production-ready implementation",
            "",
            "Impact:",
            "  • Enables efficient long-term weather prediction",
            "  • Reduces manual processing time",
            "  • Provides foundation for further research"
        ]
    )
    
    # Slide 25: Future Work
    add_content_slide(
        "Future Work",
        [
            "Model Improvements:",
            "  • Increase training epochs",
            "  • Hyperparameter optimization",
            "  • Architecture enhancements",
            "  • Additional evaluation metrics",
            "",
            "Data Enhancements:",
            "  • Higher temporal resolution",
            "  • Additional weather variables",
            "  • Longer time series",
            "",
            "System Enhancements:",
            "  • Real-time prediction",
            "  • Visualization tools",
            "  • Model interpretability",
            "  • Deployment optimization"
        ]
    )
    
    # Slide 26: Thank You
    add_title_slide(
        "Thank You",
        "Questions & Discussion"
    )
    
    return prs

if __name__ == "__main__":
    import sys
    sys.stdout.flush()
    print("Creating HMPF PowerPoint presentation...")
    sys.stdout.flush()
    try:
        presentation = create_presentation()
        output_path = r"c:\Users\91890\Downloads\HMPF_Presentation.pptx"
        presentation.save(output_path)
        print(f"Presentation saved to: {output_path}")
        print(f"Total slides: {len(presentation.slides)}")
        print("Presentation created successfully!")
        sys.stdout.flush()
    except Exception as e:
        print(f"Error creating presentation: {e}")
        import traceback
        traceback.print_exc()
        sys.stdout.flush()
